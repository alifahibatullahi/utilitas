"""Generate PowerPoint for PowerOps innovation presentation.

Theme: Light + Agriculture technology (green accent palette).
Output: d:/INOVASI 2025/PowerOps/PowerOps_Presentasi_Inovasi.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Palette (light theme, agriculture tech) ----------
BG_LIGHT = RGBColor(0xF7, 0xFA, 0xF5)   # off-white with green tint
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_DARK = RGBColor(0x1B, 0x5E, 0x20)   # deep forest green (primary)
GREEN_MAIN = RGBColor(0x2E, 0x7D, 0x32)   # main green
GREEN_LEAF = RGBColor(0x66, 0xBB, 0x6A)   # leaf green (accent)
GREEN_SOFT = RGBColor(0xE8, 0xF5, 0xE9)   # very light green (card bg)
EARTH_BROWN = RGBColor(0x8D, 0x6E, 0x63)  # soil/earth accent
SUN_GOLD = RGBColor(0xF9, 0xA8, 0x25)     # warm highlight
TEXT_DARK = RGBColor(0x22, 0x33, 0x2B)    # near-black with green hint
TEXT_MUTED = RGBColor(0x5C, 0x6B, 0x63)
DIVIDER = RGBColor(0xC8, 0xE6, 0xC9)
ALERT_RED = RGBColor(0xC6, 0x28, 0x28)
WARN_AMBER = RGBColor(0xEF, 0x6C, 0x00)

# ---------- Slide dimensions (16:9) ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line or Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_rounded(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.10
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = [text]
    else:
        lines = list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT_DARK,
                bullet_color=GREEN_MAIN, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet
        r1 = p.add_run()
        r1.text = "●  "
        r1.font.name = "Calibri"
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        # label + value (support tuple for bold prefix)
        if isinstance(item, tuple):
            prefix, rest = item
            r2 = p.add_run()
            r2.text = prefix
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.bold = True
            r3 = p.add_run()
            r3.text = rest
            r3.font.name = "Calibri"
            r3.font.size = Pt(size)
            r3.font.color.rgb = color
        else:
            r2 = p.add_run()
            r2.text = item
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
    return tb


def add_leaf_accent(slide, x, y, size=Inches(0.6)):
    """Small leaf-shaped decoration in top-right of header."""
    shape = slide.shapes.add_shape(MSO_SHAPE.TEAR, x, y, size, size)
    shape.rotation = -30
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_LEAF
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_header(slide, title, subtitle=None, slide_num=None, total=None):
    """Consistent light header with green accent bar + leaf."""
    # top accent strip
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), GREEN_MAIN)
    # leaf decorations top-right
    add_leaf_accent(slide, Inches(12.55), Inches(0.25), size=Inches(0.45))
    add_leaf_accent(slide, Inches(12.25), Inches(0.45), size=Inches(0.35))
    # title
    add_text(slide, Inches(0.5), Inches(0.30), Inches(10.5), Inches(0.7),
             title, size=28, bold=True, color=GREEN_DARK)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(10.5), Inches(0.4),
                 subtitle, size=14, color=TEXT_MUTED)
    # underline
    add_rect(slide, Inches(0.5), Inches(1.35), Inches(1.2), Pt(3.5), GREEN_LEAF)
    # slide number
    if slide_num is not None:
        add_text(slide, Inches(12.2), Inches(7.1), Inches(1.0), Inches(0.3),
                 f"{slide_num} / {total}", size=10, color=TEXT_MUTED,
                 align=PP_ALIGN.RIGHT)
    # bottom brand strip
    add_rect(slide, 0, Inches(7.38), SLIDE_W, Inches(0.12), GREEN_LEAF)
    add_text(slide, Inches(0.5), Inches(7.12), Inches(6), Inches(0.25),
             "PowerOps  •  Inovasi Digital UBB", size=10,
             color=TEXT_MUTED)


# ---------- Build presentation ----------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

TOTAL = 16


# ===== Slide 1: Cover =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_LIGHT)

# soft green gradient blocks
add_rect(s, 0, 0, SLIDE_W, Inches(2.0), GREEN_SOFT)
add_rect(s, 0, Inches(6.9), SLIDE_W, Inches(0.6), GREEN_MAIN)
add_rect(s, 0, Inches(6.75), SLIDE_W, Inches(0.15), GREEN_LEAF)

# decorative leaves
add_leaf_accent(s, Inches(12.0), Inches(0.5), size=Inches(0.9))
add_leaf_accent(s, Inches(11.4), Inches(0.9), size=Inches(0.6))
add_leaf_accent(s, Inches(0.4), Inches(6.0), size=Inches(0.7))

# centered title card
card = add_rounded(s, Inches(1.5), Inches(2.3), Inches(10.3), Inches(3.8),
                   BG_WHITE, line_color=DIVIDER)
add_text(s, Inches(1.8), Inches(2.55), Inches(9.7), Inches(0.5),
         "INOVASI DIGITAL • OPERASI UBB", size=14, bold=True,
         color=GREEN_MAIN, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.8), Inches(3.05), Inches(9.7), Inches(1.3),
         "PowerOps", size=72, bold=True, color=GREEN_DARK,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1.8), Inches(4.35), Inches(9.7), Inches(0.8),
         "Sistem Digital Terpadu untuk Operasi Utilitas Batubara",
         size=22, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.8), Inches(4.95), Inches(9.7), Inches(0.6),
         "Dari Buku, HT, & Google Sheets  →  Satu Platform Real-time",
         size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# separator
add_rect(s, Inches(5.8), Inches(5.75), Inches(1.7), Pt(2.5), GREEN_LEAF)

add_text(s, Inches(1.8), Inches(5.85), Inches(9.7), Inches(0.35),
         "Alif Ahibatullahi  •  Operator UBB  •  24 April 2026",
         size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(7.05), SLIDE_W - Inches(1), Inches(0.35),
         "Teknologi untuk Mendukung Operasi Pabrik Agro-Industri",
         size=11, color=BG_WHITE, align=PP_ALIGN.CENTER)


# ===== Slide 2: Agenda =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_LIGHT)
add_header(s, "Agenda Presentasi",
           "Alur diskusi — dari masalah menuju solusi",
           slide_num=2, total=TOTAL)

agenda = [
    ("01", "Latar Belakang", "Kondisi & masalah operasional saat ini"),
    ("02", "Gagasan: PowerOps", "Apa itu dan untuk siapa"),
    ("03", "Fitur Utama", "6 fitur menjawab 4 pain point"),
    ("04", "Workflow Baru vs Lama", "Perbandingan nyata di lapangan"),
    ("05", "Dampak & Manfaat", "Estimasi hemat waktu, kualitas data, biaya"),
    ("06", "Teknologi & Roadmap", "Stack, keamanan, rencana lanjut"),
]
x0 = Inches(0.7)
y0 = Inches(1.8)
cw = Inches(6.0)
ch = Inches(0.8)
gap_y = Inches(0.15)
for i, (num, title, desc) in enumerate(agenda):
    col = i % 2
    row = i // 2
    x = x0 + col * (cw + Inches(0.2))
    y = y0 + row * (ch + gap_y)
    card = add_rounded(s, x, y, cw, ch, BG_WHITE, line_color=DIVIDER)
    # number badge
    badge = add_rounded(s, x + Inches(0.15), y + Inches(0.15),
                        Inches(0.55), Inches(0.5), GREEN_MAIN)
    add_text(s, x + Inches(0.15), y + Inches(0.15),
             Inches(0.55), Inches(0.5), num, size=16, bold=True,
             color=BG_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.85), y + Inches(0.1),
             cw - Inches(1.0), Inches(0.4), title, size=16, bold=True,
             color=GREEN_DARK)
    add_text(s, x + Inches(0.85), y + Inches(0.42),
             cw - Inches(1.0), Inches(0.35), desc, size=11,
             color=TEXT_MUTED)

# bottom note
add_rounded(s, Inches(0.7), Inches(5.8), Inches(11.9), Inches(0.9),
            GREEN_SOFT, line_color=DIVIDER)
add_text(s, Inches(1.0), Inches(5.95), Inches(11.3), Inches(0.6),
         ["Durasi target: 15–20 menit  •  Disertai demo live  •  QnA di penutup"],
         size=13, bold=True, color=GREEN_DARK, anchor=MSO_ANCHOR.MIDDLE)


# ===== Slide 3: Latar Belakang (4 pain points) =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_LIGHT)
add_header(s, "Latar Belakang — Kondisi Saat Ini",
           "4 masalah operasional yang ingin kita selesaikan",
           slide_num=3, total=TOTAL)

pains = [
    ("01", "Tidak ada dashboard parameter aktual UBB",
     "Manajemen & CCR tidak punya single source of truth untuk status operasi real-time. Keputusan lambat."),
    ("02", "Level tanki harus ditanya via HT tiap 2 jam",
     "Operator CCR telepon Utilitas 1 & Demin 3 untuk update RCW/Demin. Jalur radio padat, tidak ada history level, tidak ada list kedatangan/permintaan solar."),
    ("03", "Laporan shift manual (buku → hitung → Sheets)",
     "Operator tulis buku + hitung kalkulator, foreman ketik ulang di Google Sheets. Rawan salah hitung, double input, tidak ada catatan fly ash."),
    ("04", "Critical & Maintenance tidak terhubung",
     "Laporan critical dan maintenance dibuat terpisah. History preventif/modifikasi per item tidak tersusun rapi, sulit cari tren kerusakan."),
]
x0 = Inches(0.6)
y0 = Inches(1.75)
cw = Inches(6.05)
ch = Inches(2.55)
for i, (num, title, desc) in enumerate(pains):
    col = i % 2
    row = i // 2
    x = x0 + col * (cw + Inches(0.2))
    y = y0 + row * (ch + Inches(0.15))
    add_rounded(s, x, y, cw, ch, BG_WHITE, line_color=DIVIDER)
    # left accent bar (red for problem)
    add_rect(s, x, y + Inches(0.15), Inches(0.08), ch - Inches(0.3),
             ALERT_RED)
    # number
    add_text(s, x + Inches(0.3), y + Inches(0.2),
             Inches(0.8), Inches(0.5), num, size=26, bold=True,
             color=ALERT_RED)
    # title
    add_text(s, x + Inches(1.1), y + Inches(0.2),
             cw - Inches(1.3), Inches(1.0),
             title, size=15, bold=True, color=GREEN_DARK)
    # desc
    add_text(s, x + Inches(0.3), y + Inches(1.15),
             cw - Inches(0.55), ch - Inches(1.3),
             desc, size=12, color=TEXT_DARK)


# ===== Slide 4: Gagasan =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_LIGHT)
add_header(s, "Gagasan: Apa itu PowerOps?",
           "Satu aplikasi web terpadu untuk seluruh workflow operasi UBB",
           slide_num=4, total=TOTAL)

# big tagline card
add_rounded(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(1.4),
            GREEN_SOFT, line_color=GREEN_LEAF)
add_text(s, Inches(0.6), Inches(1.85), Inches(12.1), Inches(0.5),
         "“Satu Layar untuk Satu Operasi”",
         size=32, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(2.55), Inches(12.1), Inches(0.5),
         "Menggantikan buku shift, panggilan HT, dan Google Sheets yang tersebar — dalam satu platform real-time",
         size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# 4 feature pillars
pillars = [
    ("\U0001F4F1", "Multi-device", "Desktop CCR, HP operator, HP foreman, HP manajemen"),
    ("⚡", "Real-time", "Data update otomatis via Supabase subscription — tanpa refresh"),
    ("✓", "Sudah Berjalan", "Bukan konsep. Sudah dipakai, tinggal scale-up adopsi"),
    ("\U0001F4B0", "Biaya Rendah", "In-house, stack open source, layanan tier gratis/murah"),
]
x0 = Inches(0.6)
y0 = Inches(3.45)
cw = Inches(2.95)
ch = Inches(3.3)
for i, (icon, title, desc) in enumerate(pillars):
    x = x0 + i * (cw + Inches(0.12))
    add_rounded(s, x, y0, cw, ch, BG_WHITE, line_color=DIVIDER)
    # icon circle
    add_rounded(s, x + Inches(1.1), y0 + Inches(0.3),
                Inches(0.75), Inches(0.75), GREEN_MAIN)
    add_text(s, x + Inches(1.1), y0 + Inches(0.3),
             Inches(0.75), Inches(0.75), icon, size=24, bold=True,
             color=BG_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), y0 + Inches(1.25),
             cw - Inches(0.3), Inches(0.5), title, size=16, bold=True,
             color=GREEN_DARK, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y0 + Inches(1.8),
             cw - Inches(0.4), ch - Inches(1.8), desc, size=12,
             color=TEXT_DARK, align=PP_ALIGN.CENTER)


# ===== Slide 5: Arsitektur Solusi =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_LIGHT)
add_header(s, "Arsitektur Solusi",
           "3 lapis sederhana: sumber data → platform → konsumen",
           slide_num=5, total=TOTAL)

# 3 horizontal bands
def band(y, title, desc_items, color, icon):
    add_rounded(s, Inches(0.7), y, Inches(11.9), Inches(1.55), BG_WHITE,
                line_color=DIVIDER)
    add_rect(s, Inches(0.7), y, Inches(0.12), Inches(1.55), color)
    # icon badge
    add_rounded(s, Inches(1.0), y + Inches(0.3),
                Inches(0.9), Inches(0.9), color)
    add_text(s, Inches(1.0), y + Inches(0.3), Inches(0.9), Inches(0.9),
             icon, size=28, bold=True, color=BG_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title
    add_text(s, Inches(2.1), y + Inches(0.25), Inches(10.2), Inches(0.45),
             title, size=17, bold=True, color=GREEN_DARK)
    # desc pills
    px = Inches(2.1)
    for item in desc_items:
        est_w = Inches(0.35 + 0.095 * len(item))
        pill = add_rounded(s, px, y + Inches(0.85),
                           est_w, Inches(0.5), GREEN_SOFT,
                           line_color=GREEN_LEAF)
        add_text(s, px, y + Inches(0.85), est_w, Inches(0.5),
                 item, size=11, bold=True, color=GREEN_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        px = px + est_w + Inches(0.15)


band(Inches(1.75), "SUMBER DATA",
     ["Operator (input shift)", "Foreman (critical/maintenance)",
      "Meter/Sensor (manual entry)"],
     GREEN_MAIN, "\U0001F4E5")

band(Inches(3.45), "PLATFORM PowerOps",
     ["Dashboard", "Tank Level", "Logsheet", "Critical+Maintenance",
      "History", "Kanban"],
     GREEN_DARK, "▦")

band(Inches(5.15), "KONSUMEN DATA",
     ["Manajemen (dashboard)", "CCR (tank level)",
      "Auto-sync Google Sheets (laporan lama)"],
     EARTH_BROWN, "\U0001F4E4")

# arrows between bands
for y in (Inches(3.25), Inches(4.95)):
    arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                             Inches(6.55), y, Inches(0.25), Inches(0.22))
    arr.fill.solid()
    arr.fill.fore_color.rgb = GREEN_LEAF
    arr.line.fill.background()

# bottom tag
add_text(s, Inches(0.7), Inches(6.95), Inches(11.9), Inches(0.35),
         "Kompatibilitas: laporan Google Sheets lama tetap berjalan paralel — peralihan bertahap, nol disrupsi",
         size=12, bold=True, color=EARTH_BROWN, align=PP_ALIGN.CENTER)


# ===== Feature slide builder =====
def feature_slide(num, feature_no, title_short, page_url, bullets, solves):
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header(s, f"Fitur {feature_no} — {title_short}",
               f"Halaman: {page_url}",
               slide_num=num, total=TOTAL)

    # screenshot placeholder (left)
    ph_x, ph_y = Inches(0.6), Inches(1.75)
    ph_w, ph_h = Inches(6.5), Inches(4.5)
    add_rounded(s, ph_x, ph_y, ph_w, ph_h, BG_WHITE, line_color=GREEN_LEAF)
    # browser bar
    add_rect(s, ph_x + Inches(0.15), ph_y + Inches(0.15),
             ph_w - Inches(0.3), Inches(0.35), GREEN_SOFT)
    for i, col in enumerate([ALERT_RED, WARN_AMBER, GREEN_LEAF]):
        add_rounded(s, ph_x + Inches(0.28) + i * Inches(0.22),
                    ph_y + Inches(0.22), Inches(0.17), Inches(0.17), col)
    add_text(s, ph_x + Inches(1.0), ph_y + Inches(0.19),
             ph_w - Inches(1.2), Inches(0.3),
             f"powerops.local{page_url}", size=10, color=TEXT_MUTED)
    # placeholder center text
    add_text(s, ph_x, ph_y + Inches(1.8), ph_w, Inches(0.5),
             "[ Screenshot halaman ]", size=16, bold=True,
             color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_text(s, ph_x, ph_y + Inches(2.3), ph_w, Inches(0.4),
             "Ganti area ini dengan screenshot aktual saat finalisasi",
             size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # right: bullets
    bx, by = Inches(7.3), Inches(1.75)
    bw, bh = Inches(5.4), Inches(4.5)
    add_rounded(s, bx, by, bw, bh, BG_WHITE, line_color=DIVIDER)
    add_text(s, bx + Inches(0.3), by + Inches(0.2), bw - Inches(0.6),
             Inches(0.5), "Kemampuan Utama", size=16, bold=True,
             color=GREEN_DARK)
    add_rect(s, bx + Inches(0.3), by + Inches(0.72),
             Inches(0.6), Pt(2.5), GREEN_LEAF)
    add_bullets(s, bx + Inches(0.3), by + Inches(0.9),
                bw - Inches(0.6), bh - Inches(1.1), bullets,
                size=13, line_spacing=1.25)

    # solves badge
    add_rounded(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.7),
                GREEN_MAIN)
    add_text(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.7),
             f"✓  {solves}", size=14, bold=True, color=BG_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ===== Slide 6: Fitur 1 Dashboard =====
feature_slide(
    6, "1", "Dashboard Parameter Aktual UBB", "/dashboard",
    [
        ("Boiler A & B: ", "steam flow, furnace temp, vakum, O2, coal feeder (3 feeder)"),
        ("STG: ", "load MW, frequency, steam inlet/condensate, vakum, thrust bearing"),
        ("Distribusi steam: ", "pabrik 1, pabrik 3, inlet turbin"),
        ("Quick stats: ", "total steam, total coal, trending indicator"),
        ("Kartu visual: ", "angka besar, gauge, status warna"),
        ("Real-time: ", "update otomatis tanpa refresh manual"),
    ],
    "Menjawab Masalah #1: Manajemen & CCR kini punya single source of truth"
)


# ===== Slide 7: Fitur 2 Tank Level =====
feature_slide(
    7, "2", "Monitoring Level Tanki (CCR)", "/tank-level",
    [
        ("Visual glass tank: ", "DEMIN (1.250 m³), RCW (4.600 m³), SOLAR (2×200 m³)"),
        ("Status warna: ", "Normal (hijau), Warning (kuning), Kritis (merah)"),
        ("Flow in/out real-time: ", "t/h ke/dari masing-masing sumber/tujuan"),
        ("Trend history: ", "24 data point terakhir dalam line chart"),
        ("Tracking pompa: ", "“Aktif sejak [tanggal jam]” untuk Demin Revamp"),
        ("Subscription Supabase: ", "level berubah → UI otomatis update"),
    ],
    "Menjawab Masalah #2: CCR tidak perlu HT lagi tiap 2 jam"
)


# ===== Slide 8: Fitur 3 Solar =====
feature_slide(
    8, "3", "Solar Unloading & Permintaan", "/tank-level (modal Solar)",
    [
        ("Log kedatangan: ", "tanggal, liter, supplier/perusahaan"),
        ("Log permintaan: ", "tanggal, liter, tujuan (mis. Boiler)"),
        ("Tabs terpisah: ", "unloading vs permintaan dengan paginasi 5/halaman"),
        ("CRUD lengkap: ", "edit & hapus dengan konfirmasi"),
        ("Terintegrasi tanki SOLAR: ", "level & history di satu layar"),
        ("Audit trail: ", "semua perubahan tercatat rapi"),
    ],
    "Menjawab Masalah #2 (lanjut): Nota fisik & HT request hilang"
)


# ===== Slide 9: Fitur 4 Logsheet =====
feature_slide(
    9, "4", "Logsheet Digital per Shift", "/input-shift",
    [
        ("3 shift: ", "Malam 06:00  •  Pagi 14:00  •  Sore 22:00"),
        ("9 tab: ", "Boiler A/B, Turbin, Generator, Distribusi Steam, Handling, ESP, Coal Bunker, Lab/QC"),
        ("40+ parameter: ", "tercatat terstruktur per shift"),
        ("Auto-calculation: ", "CR, selisih totalizer, konsumsi batubara/steam"),
        ("Unloading fly ash (tab ESP): ", "silo A/B, supplier, tujuan, ritase"),
        ("Auto-sync Google Sheets: ", "laporan lama tetap dapat diakses"),
    ],
    "Menjawab Masalah #3: Tidak ada lagi double input & hitung manual"
)


# ===== Slide 10: Fitur 5 LHUBB =====
feature_slide(
    10, "5", "Laporan Harian (LHUBB) Otomatis", "/laporan-harian",
    [
        ("Agregasi 24 jam: ", "steam, power, coal, unloading fly ash A/B"),
        ("Foreman tinggal review: ", "tidak perlu ketik ulang dari buku operator"),
        ("Preview siap cetak: ", "layout laporan resmi"),
        ("Kirim via WhatsApp: ", "integrasi notifikasi otomatis"),
        ("Data turunan otomatis: ", "dari tabel input-shift harian"),
        ("History tersimpan: ", "semua laporan harian bisa ditelusuri"),
    ],
    "Menjawab Masalah #3 (lanjut): Foreman fokus review, bukan entry ulang"
)


# ===== Slide 11: Fitur 6 Critical + Maintenance =====
feature_slide(
    11, "6", "Critical + Maintenance Terintegrasi", "/critical",
    [
        ("Daftar critical equipment: ", "item, deskripsi, scope, foreman, status, notif"),
        ("Maintenance log per item: ", "corrective, preventif, modifikasi"),
        ("Timeline activity: ", "siapa mengubah apa, kapan"),
        ("Filter per equipment: ", "lihat semua history dalam satu layar"),
        ("Status auto-sync: ", "semua maintenance OK → critical bisa close"),
        ("Work orders: ", "perencanaan preventif/modifikasi terjadwal"),
    ],
    "Menjawab Masalah #4: History per equipment kini tertata & tertelusuri"
)


# ===== Slide 12: Workflow Sebelum vs Sesudah =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_LIGHT)
add_header(s, "Workflow: Sebelum vs Sesudah",
           "Perbandingan kongkret proses operasional harian",
           slide_num=12, total=TOTAL)

# Table header
tx = Inches(0.6)
ty = Inches(1.75)
col_w = [Inches(3.3), Inches(4.4), Inches(4.6)]
row_h = Inches(0.7)
headers = ["Proses", "Sebelum", "Sesudah dengan PowerOps"]

# header row
x = tx
for i, hdr in enumerate(headers):
    add_rect(s, x, ty, col_w[i], Inches(0.55), GREEN_DARK)
    add_text(s, x, ty, col_w[i], Inches(0.55), hdr, size=14, bold=True,
             color=BG_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x = x + col_w[i]

rows = [
    ("Update level tanki ke CCR",
     "Telepon HT tiap 2 jam",
     "Buka layar, data auto-refresh"),
    ("Laporan shift",
     "Tulis buku → hitung kalkulator → foreman ketik ulang Sheets",
     "Operator input di app (autocalculate) → foreman review & finalisasi"),
    ("Cek history parameter",
     "Buka buku shift fisik lama",
     "Halaman History dengan filter & custom column"),
    ("Catat critical",
     "Form terpisah, arsip manual",
     "1 entry → link ke maintenance → timeline tersimpan"),
    ("Cek level & history solar",
     "Tanya handling / lihat nota fisik",
     "Modal solar dengan list unloading & permintaan"),
    ("Cari history preventif pompa X",
     "Bongkar arsip per bulan",
     "Filter item “Pompa X” → semua record muncul"),
]

ry = ty + Inches(0.55)
for i, row in enumerate(rows):
    x = tx
    bg = BG_WHITE if i % 2 == 0 else GREEN_SOFT
    for j, val in enumerate(row):
        add_rect(s, x, ry, col_w[j], Inches(0.65), bg,
                 line=Pt(0.5), line_color=DIVIDER)
        color = GREEN_DARK if j == 0 else (ALERT_RED if j == 1 else GREEN_MAIN)
        bold = (j == 0)
        size = 11 if j > 0 else 12
        add_text(s, x + Inches(0.1), ry, col_w[j] - Inches(0.2),
                 Inches(0.65), val, size=size, bold=bold, color=color,
                 anchor=MSO_ANCHOR.MIDDLE)
        x = x + col_w[j]
    ry = ry + Inches(0.65)

# legend
add_text(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3),
         "Kolom “Sebelum” (merah) = proses manual saat ini  •  Kolom “Sesudah” (hijau) = proses dengan PowerOps",
         size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ===== Slide 13: Dampak =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_LIGHT)
add_header(s, "Dampak & Manfaat",
           "Estimasi konservatif — validasi dengan data aktual lapangan",
           slide_num=13, total=TOTAL)

# Big metric cards (top row)
metrics = [
    ("~30", "menit/shift", "hemat dari tidak ketik ulang laporan", GREEN_MAIN),
    ("~1,5", "jam/hari", "akumulasi 3 shift per unit", GREEN_DARK),
    ("~12", "panggilan HT/hari", "hilang — jalur radio lebih senggang", EARTH_BROWN),
    ("0", "biaya lisensi", "stack open source, tier gratis/murah", SUN_GOLD),
]
x0 = Inches(0.6)
y0 = Inches(1.75)
cw = Inches(2.97)
ch = Inches(2.0)
for i, (big, unit, desc, color) in enumerate(metrics):
    x = x0 + i * (cw + Inches(0.12))
    add_rounded(s, x, y0, cw, ch, BG_WHITE, line_color=DIVIDER)
    # top bar
    add_rect(s, x, y0, cw, Inches(0.15), color)
    add_text(s, x + Inches(0.15), y0 + Inches(0.3), cw - Inches(0.3),
             Inches(0.85), big, size=44, bold=True, color=color,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y0 + Inches(1.15), cw - Inches(0.3),
             Inches(0.35), unit, size=13, bold=True, color=GREEN_DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y0 + Inches(1.5), cw - Inches(0.3),
             Inches(0.45), desc, size=11, color=TEXT_MUTED,
             align=PP_ALIGN.CENTER)

# Bottom row: 4 benefit categories
benefits = [
    ("Kualitas Data",
     ["Auto-calc hilangkan salah hitung manual",
      "100% history digital → audit cepat",
      "Single source of truth — tidak ada versi bentrok"]),
    ("Respons Operasional",
     ["Manajemen & CCR lihat realtime kapan saja",
      "Abnormal tanki terdeteksi via status warna",
      "Notifikasi WhatsApp untuk kondisi kritis"]),
    ("Tata Kelola Maintenance",
     ["History preventif/modifikasi per equipment",
      "Dasar analisis keandalan & root cause",
      "Critical tidak lagi “lupa ditutup”"]),
]
x0 = Inches(0.6)
y0 = Inches(3.95)
cw = Inches(4.03)
ch = Inches(3.0)
for i, (title, items) in enumerate(benefits):
    x = x0 + i * (cw + Inches(0.12))
    add_rounded(s, x, y0, cw, ch, GREEN_SOFT, line_color=GREEN_LEAF)
    add_text(s, x + Inches(0.25), y0 + Inches(0.15), cw - Inches(0.5),
             Inches(0.5), title, size=15, bold=True, color=GREEN_DARK)
    add_rect(s, x + Inches(0.25), y0 + Inches(0.65),
             Inches(0.5), Pt(2.5), GREEN_MAIN)
    add_bullets(s, x + Inches(0.25), y0 + Inches(0.8),
                cw - Inches(0.5), ch - Inches(0.95), items,
                size=11, line_spacing=1.2)


# ===== Slide 14: Teknologi =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_LIGHT)
add_header(s, "Teknologi yang Dipakai",
           "Stack modern, standar industri, mudah dimaintain",
           slide_num=14, total=TOTAL)

tech_groups = [
    ("Frontend", [
        "Next.js 16 (React 19)",
        "TypeScript 5",
        "Tailwind CSS 4",
        "Recharts (grafik)",
    ], GREEN_MAIN),
    ("Backend & Database", [
        "Supabase (PostgreSQL)",
        "Real-time subscription",
        "Role-based access control",
        "Server-Side Rendering",
    ], GREEN_DARK),
    ("Integrasi", [
        "Google Sheets API (sync)",
        "AWS S3 (file lampiran)",
        "WhatsApp (notifikasi)",
        "Auth operator (grup A/B/C/D)",
    ], EARTH_BROWN),
    ("Akses & Deploy", [
        "Web browser apa saja",
        "Desktop CCR, HP mobile",
        "Hosting Vercel",
        "Role: op/foreman/supervisor/admin",
    ], SUN_GOLD),
]
x0 = Inches(0.6)
y0 = Inches(1.75)
cw = Inches(2.97)
ch = Inches(4.2)
for i, (grp, items, color) in enumerate(tech_groups):
    x = x0 + i * (cw + Inches(0.12))
    add_rounded(s, x, y0, cw, ch, BG_WHITE, line_color=DIVIDER)
    add_rect(s, x, y0, cw, Inches(0.55), color)
    add_text(s, x, y0, cw, Inches(0.55), grp, size=15, bold=True,
             color=BG_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, x + Inches(0.2), y0 + Inches(0.75),
                cw - Inches(0.4), ch - Inches(0.9), items,
                size=12, line_spacing=1.35, bullet_color=color)

# bottom tag
add_rounded(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.75),
            GREEN_SOFT, line_color=GREEN_LEAF)
add_text(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.75),
         "Source code & dokumentasi di Git • Stack standar industri • Knowledge transfer siap dilakukan",
         size=13, bold=True, color=GREEN_DARK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ===== Slide 15: Roadmap & Penutup =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_LIGHT)
add_header(s, "Roadmap Singkat & Penutup",
           "Yang sudah jalan, sedang dikerjakan, dan rencana lanjut",
           slide_num=15, total=TOTAL)

# 3 status columns
cols = [
    ("SUDAH JALAN", GREEN_MAIN, "✓", [
        "Dashboard parameter aktual",
        "Tank Level (DEMIN/RCW/SOLAR)",
        "Solar unloading & permintaan",
        "Logsheet 3 shift × 9 tab",
        "Laporan Harian (LHUBB)",
        "Critical + Maintenance",
        "History viewer",
        "Kanban task board",
        "Sync Google Sheets",
    ]),
    ("SEDANG DISEMPURNAKAN", SUN_GOLD, "⊛", [
        "Dashboard stream days boiler (hari sejak flow > 0)",
        "Notifikasi WhatsApp otomatis untuk kondisi kritis",
        "Penyempurnaan UX laporan",
    ]),
    ("RENCANA LANJUT", EARTH_BROWN, "▶", [
        "Prediksi/analitik tren kerusakan",
        "Integrasi langsung sensor (OPC UA)",
        "Mobile app native",
    ]),
]
x0 = Inches(0.6)
y0 = Inches(1.75)
cw = Inches(4.03)
ch = Inches(3.7)
for i, (title, color, icon, items) in enumerate(cols):
    x = x0 + i * (cw + Inches(0.12))
    add_rounded(s, x, y0, cw, ch, BG_WHITE, line_color=DIVIDER)
    add_rect(s, x, y0, cw, Inches(0.6), color)
    add_text(s, x, y0, cw, Inches(0.6),
             f"{icon}  {title}", size=14, bold=True, color=BG_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, x + Inches(0.2), y0 + Inches(0.8),
                cw - Inches(0.4), ch - Inches(1.0), items,
                size=11, line_spacing=1.3, bullet_color=color)

# Ask for support
add_rounded(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.4),
            GREEN_DARK)
add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.45),
         "Dukungan yang Dibutuhkan dari Manajemen",
         size=16, bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.8),
         ["(1) Legitimasi PowerOps sebagai tool resmi UBB   •   "
          "(2) Standarisasi pemakaian lintas shift   •   "
          "(3) Bantuan onboarding seluruh operator & foreman"],
         size=13, color=BG_WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)


# ===== Slide 16: Thank You & QnA =====
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_LIGHT)

# full decoration
add_rect(s, 0, 0, SLIDE_W, Inches(3.0), GREEN_SOFT)
add_leaf_accent(s, Inches(0.6), Inches(0.6), size=Inches(1.1))
add_leaf_accent(s, Inches(1.6), Inches(0.4), size=Inches(0.7))
add_leaf_accent(s, Inches(11.6), Inches(1.7), size=Inches(1.0))
add_leaf_accent(s, Inches(12.1), Inches(2.2), size=Inches(0.6))
add_rect(s, 0, Inches(6.9), SLIDE_W, Inches(0.6), GREEN_MAIN)
add_rect(s, 0, Inches(6.75), SLIDE_W, Inches(0.15), GREEN_LEAF)

# main text
add_text(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.3),
         "Terima Kasih", size=72, bold=True, color=GREEN_DARK,
         align=PP_ALIGN.CENTER)
add_rect(s, Inches(6.0), Inches(3.55), Inches(1.3), Pt(3.5), GREEN_LEAF)
add_text(s, Inches(0.5), Inches(3.75), Inches(12.3), Inches(0.6),
         "Pertanyaan & Diskusi", size=28, bold=True, color=TEXT_DARK,
         align=PP_ALIGN.CENTER)

# Contact card
add_rounded(s, Inches(3.5), Inches(4.75), Inches(6.3), Inches(1.6),
            BG_WHITE, line_color=DIVIDER)
add_text(s, Inches(3.5), Inches(4.9), Inches(6.3), Inches(0.4),
         "Alif Ahibatullahi", size=18, bold=True, color=GREEN_DARK,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(3.5), Inches(5.3), Inches(6.3), Inches(0.4),
         "Operator UBB  •  PowerOps Developer",
         size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(3.5), Inches(5.75), Inches(6.3), Inches(0.4),
         "alifahibatullahi@gmail.com",
         size=12, color=GREEN_MAIN, align=PP_ALIGN.CENTER)

# footer
add_text(s, Inches(0.5), Inches(7.05), SLIDE_W - Inches(1), Inches(0.35),
         "PowerOps  •  Teknologi Digital untuk Operasi Pabrik Agro-Industri",
         size=11, color=BG_WHITE, align=PP_ALIGN.CENTER)


# ---------- Save ----------
out_path = r"d:\INOVASI 2025\PowerOps\PowerOps_Presentasi_Inovasi.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
