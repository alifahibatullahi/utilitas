#!/usr/bin/env python3
"""Build PowerOps presentation with agroindustri theme."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ----- Theme: Agroindustri -----
GREEN_DEEP   = RGBColor(0x1B, 0x5E, 0x20)   # primary
GREEN_MAIN   = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LIGHT  = RGBColor(0xA5, 0xD6, 0xA7)
GREEN_BG     = RGBColor(0xF1, 0xF8, 0xF4)   # slide background
AMBER        = RGBColor(0xFF, 0xA0, 0x00)   # harvest accent
BROWN        = RGBColor(0x6D, 0x4C, 0x41)   # earth tone
INK          = RGBColor(0x1B, 0x30, 0x24)   # text dark
INK_SOFT     = RGBColor(0x5C, 0x7A, 0x6A)   # text muted
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
RED          = RGBColor(0xC6, 0x28, 0x28)
BLUE         = RGBColor(0x15, 0x65, 0xC0)

HERE = os.path.dirname(os.path.abspath(__file__))

# ----- Setup 16:9 deck -----
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=GREEN_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_header_bar(slide, height=Inches(0.55)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, height)
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN_DEEP
    bar.line.fill.background()
    # decorative amber stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, height, SW, Inches(0.06))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = AMBER
    stripe.line.fill.background()
    return bar


def add_footer(slide, page_no=None, total=None):
    foot = slide.shapes.add_textbox(Inches(0.4), SH - Inches(0.4),
                                    SW - Inches(0.8), Inches(0.3))
    tf = foot.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "PowerOps · Digitalisasi Operasi UBB · Internal Use"
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.color.rgb = INK_SOFT
    if page_no is not None and total is not None:
        rp = p.add_run()
        rp.text = f"          {page_no} / {total}"
        rp.font.size = Pt(10)
        rp.font.color.rgb = INK_SOFT


def add_textbox(slide, left, top, width, height, text, *,
                size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def set_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_image_fit(slide, path, left, top, max_w, max_h):
    """Insert image fitting within (max_w, max_h) preserving aspect ratio."""
    from PIL import Image
    im = Image.open(path)
    iw, ih = im.size
    iw_emu = Emu(iw * 9525)   # 1 px = 9525 EMU at 96 DPI baseline
    ih_emu = Emu(ih * 9525)
    # scale to fit
    sx = max_w / iw_emu
    sy = max_h / ih_emu
    s = min(sx, sy)
    w = int(iw_emu * s); h = int(ih_emu * s)
    # center
    x = left + (max_w - w) // 2
    y = top + (max_h - h) // 2
    slide.shapes.add_picture(path, x, y, width=w, height=h)


# =========================================================
# SLIDE 1 — COVER
# =========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, GREEN_BG)
# left green panel
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.2), SH)
panel.fill.solid(); panel.fill.fore_color.rgb = GREEN_DEEP
panel.line.fill.background()
# amber accent vertical
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), 0, Inches(0.08), SH)
acc.fill.solid(); acc.fill.fore_color.rgb = AMBER
acc.line.fill.background()
# leaf-like decorative circle
deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(-1.5),
                          Inches(3.5), Inches(3.5))
deco.fill.solid(); deco.fill.fore_color.rgb = GREEN_MAIN
deco.line.fill.background()

# left panel text (eyebrow + brand)
add_textbox(s, Inches(0.6), Inches(0.6), Inches(4.4), Inches(0.4),
            "PROYEK DIGITALISASI · 2026", size=12, bold=True, color=AMBER)
add_textbox(s, Inches(0.6), Inches(5.9), Inches(4.4), Inches(0.4),
            "Unit Boiler & Biomassa Agroindustri", size=14, color=GREEN_LIGHT)
add_textbox(s, Inches(0.6), Inches(6.3), Inches(4.4), Inches(0.4),
            "Presenter: Tim Operasi UBB", size=12, color=WHITE)

# main title (right side)
add_textbox(s, Inches(5.8), Inches(1.6), Inches(7.0), Inches(1.4),
            "PowerOps", size=72, bold=True, color=GREEN_DEEP)
add_textbox(s, Inches(5.8), Inches(2.9), Inches(7.0), Inches(1.0),
            "Digitalisasi Workflow Operasi UBB",
            size=28, bold=True, color=INK)
add_textbox(s, Inches(5.8), Inches(3.9), Inches(7.0), Inches(2.0),
            ["Workflow Laporan Shift & Harian",
             "Pengelompokan Critical Maintenance",
             "Dokumentasi Foto Lapangan"],
            size=18, color=INK_SOFT)

# accent line near bottom
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(6.2),
                         Inches(2.0), Inches(0.06))
line.fill.solid(); line.fill.fore_color.rgb = AMBER
line.line.fill.background()
add_textbox(s, Inches(5.8), Inches(6.35), Inches(7.0), Inches(0.4),
            "Operator fokus operasi · Supervisor fokus keputusan",
            size=13, color=INK_SOFT)

set_notes(s, """SPEAKER NOTES — SLIDE COVER

Pembuka:
- Selamat pagi/siang Bapak/Ibu. Hari ini saya ingin memaparkan inisiatif digitalisasi operasi UBB yang kami sebut PowerOps.
- Fokus presentasi ini adalah dua perbaikan utama yang sudah berjalan: (1) penyederhanaan workflow laporan shift & harian, dan (2) restrukturisasi pencatatan Critical Maintenance dengan dokumentasi foto.
- Tujuan akhirnya satu: Operator bisa fokus mengoperasikan plant, Supervisor fokus mengambil keputusan, tanpa terbebani administrasi berulang.
- Estimasi durasi: ~10-15 menit, dilanjutkan diskusi.

Konteks bisnis (untuk Q&A):
- UBB adalah backbone operasi. Setiap menit downtime dan setiap miskomunikasi berdampak ke produksi.
- PowerOps tidak menambah sistem baru yang rumit — justru menggantikan beberapa tools manual menjadi satu platform.""")

# =========================================================
# SLIDE 2 — AGENDA
# =========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s)
add_textbox(s, Inches(0.4), Inches(0.08), Inches(8), Inches(0.4),
            "Agenda", size=20, bold=True, color=WHITE)

add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
            "Apa yang akan dibahas", size=32, bold=True, color=GREEN_DEEP)
add_textbox(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5),
            "Tiga perbaikan utama hasil implementasi PowerOps di UBB",
            size=16, color=INK_SOFT)

items = [
    ("01", "Workflow Laporan Shift & Harian",
     "Dari 5 tools terpisah menjadi 1 platform terpusat. Distribusi laporan otomatis ke Washift & WhatsApp Group."),
    ("02", "Pengelompokan Critical Maintenance",
     "Critical & Maintenance dulu berdiri sendiri-sendiri. Sekarang satu Critical Job menaungi banyak Maintenance terkait — kronologi pekerjaan jadi jelas."),
    ("03", "Dokumentasi Foto Lapangan",
     "Foto wajib menyertai setiap entry Critical, Preventif, dan Modifikasi — bukti visual langsung dari lapangan, audit-ready."),
]
y = Inches(2.7)
for i, (num, title, desc) in enumerate(items):
    # number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.9), Inches(0.9))
    circ.fill.solid(); circ.fill.fore_color.rgb = GREEN_DEEP
    circ.line.fill.background()
    add_textbox(s, Inches(0.6), y + Inches(0.13), Inches(0.9), Inches(0.7),
                num, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.7), y - Inches(0.05), Inches(11), Inches(0.45),
                title, size=20, bold=True, color=INK)
    add_textbox(s, Inches(1.7), y + Inches(0.42), Inches(11), Inches(0.7),
                desc, size=14, color=INK_SOFT)
    y += Inches(1.35)

add_footer(s, 2, 6)
set_notes(s, """SPEAKER NOTES — AGENDA

- Saya akan bahas tiga hal secara berurutan.
- Yang pertama, alur laporan shift/harian — ini perbaikan paling kelihatan karena dipakai setiap hari, setiap shift.
- Yang kedua, perubahan struktur Critical Maintenance. Dulu setiap entry maintenance berdiri sendiri, jadi sulit melihat kronologi satu pekerjaan utuh. Sekarang dikelompokkan di bawah satu Critical Job.
- Yang ketiga, fitur dokumentasi foto. Ini berlaku tidak hanya untuk Critical, tapi juga Preventif dan Modifikasi.
- Setelah itu saya tutup dengan ringkasan manfaat.""")

# =========================================================
# SLIDE 3 — Workflow Laporan Shift & Harian
# =========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s)
add_textbox(s, Inches(0.4), Inches(0.08), Inches(10), Inches(0.4),
            "01 · Workflow Laporan Shift & Harian",
            size=20, bold=True, color=WHITE)

add_textbox(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
            "Dari banyak tools manual → satu platform terpusat",
            size=26, bold=True, color=GREEN_DEEP)
add_textbox(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
            "SEBELUM: operator lapor via HT → CCR catat logbook → Foreman tulis ulang di Spreadsheet → Supervisor telepon Washift.  SESUDAH: operator input langsung di PowerOps, distribusi otomatis.",
            size=12, color=INK_SOFT)

# insert flowchart image
img_path = os.path.join(HERE, "powerops-shift-report-v1.png")
add_image_fit(s, img_path,
              Inches(0.6), Inches(2.0),
              Inches(12.1), Inches(5.0))

add_footer(s, 3, 6)
set_notes(s, """SPEAKER NOTES — Workflow Laporan Shift & Harian

Penjelasan flowchart:
- Sisi kiri (SEBELUM): operator harus menghubungi CCR via HT satu per satu. CCR mencatat di logbook (tulis tangan). Lalu Foreman menulis ulang ke Spreadsheet → ini DOUBLE INPUT, sumber utama inefisiensi & risiko salah ketik. Supervisor lalu masih telepon Washift secara manual.
- Sisi kanan (SESUDAH): operator langsung input di form PowerOps. Data otomatis masuk ke database sebagai SINGLE SOURCE OF TRUTH. Foreman/Supervisor cukup ACC. Setelah ACC, sistem otomatis push laporan ke Washift & WhatsApp Group.

Angka kunci yang bisa disebut:
- 5 tools (HT, Logbook, Spreadsheet, Telepon, WA) → 1 tool (PowerOps)
- 3 orang input data (Operator + CCR + Foreman) → 1 orang (Operator saja)
- Distribusi dari MANUAL menjadi OTOMATIS

Talking point:
- "Yang paling penting di sini: operator tidak lagi 'kerja administrasi'. Mereka input sekali, dan datanya mengalir otomatis ke semua pihak yang butuh."
- "CCR yang dulu sibuk mencatat logbook sekarang bisa fokus monitoring."
- "Supervisor punya jejak digital — tidak perlu lagi 'menelepon ke Washift', semua tercatat."

Antisipasi pertanyaan:
- "Bagaimana kalau jaringan down?" → form bisa di-cache lokal, sync saat online (kalau memang sudah ada).
- "Bagaimana audit trail?" → semua entry tersimpan dengan timestamp, operator, dan jejak approval.""")

# =========================================================
# SLIDE 4 — Critical Maintenance
# =========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s)
add_textbox(s, Inches(0.4), Inches(0.08), Inches(10), Inches(0.4),
            "02 · Pengelompokan Critical Maintenance",
            size=20, bold=True, color=WHITE)

add_textbox(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
            "Satu Critical Job → banyak Maintenance terkait",
            size=26, bold=True, color=GREEN_DEEP)
add_textbox(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
            "SEBELUM: tiap maintenance dicatat sebagai entry independen — sulit lihat kronologi.  SESUDAH: dikelompokkan hierarkis di bawah satu Critical Job, alur kerja terbaca dari awal hingga selesai.",
            size=12, color=INK_SOFT)

img_path = os.path.join(HERE, "powerops-critical-maintenance.png")
add_image_fit(s, img_path,
              Inches(0.6), Inches(2.0),
              Inches(12.1), Inches(5.0))

add_footer(s, 4, 6)
set_notes(s, """SPEAKER NOTES — Critical Maintenance

Penjelasan struktur:
- SEBELUM: kalau ada masalah Pompa A misalnya, setiap tindakan (deteksi bocor, ganti seal, test running, serah terima) dicatat sebagai entry Critical TERSENDIRI. Jadi di sistem kita lihat Critical #001, #002, #003, #004 — semuanya independen, tidak saling tahu.
  - Masalahnya: kalau atasan atau auditor mau melihat "Bagaimana sih kronologi pekerjaan perbaikan Pompa A?" — kita harus menelusuri satu per satu, manual. Sulit melihat gambaran utuh.

- SESUDAH: kita perkenalkan konsep CRITICAL JOB sebagai parent/induk. Satu Critical Job ("Perbaikan Pompa A") menaungi beberapa Maintenance items (inspeksi, ganti seal, test running, closeout) secara berurutan.
  - Manfaat: kronologi terbaca dari atas ke bawah. Status pekerjaan utuh terlihat dalam satu tampilan. Mau review pekerjaan minggu lalu? Tinggal buka Critical Job-nya.

Talking point:
- "Ini bukan hanya soal rapi — ini soal organizational memory. 6 bulan lagi kalau pompa yang sama bermasalah, kita tinggal buka Critical Job lama untuk lihat apa saja yang sudah pernah dilakukan dan apa hasilnya."
- "Konsep yang sama bisa diterapkan ke Preventif & Modifikasi — semua punya struktur parent-child yang sama."

Antisipasi pertanyaan:
- "Bagaimana kalau pekerjaan berkembang di tengah jalan, butuh tambahan maintenance?" → tambah saja item baru di bawah Critical Job yang sudah ada, urutan tetap terjaga.
- "Bagaimana kalau ada Critical Job paralel?" → bisa, sistem mendukung banyak Critical Job berjalan bersamaan dengan owner masing-masing.""")

# =========================================================
# SLIDE 5 — Photo Documentation
# =========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s)
add_textbox(s, Inches(0.4), Inches(0.08), Inches(10), Inches(0.4),
            "03 · Dokumentasi Foto Lapangan",
            size=20, bold=True, color=WHITE)

add_textbox(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
            "Setiap entry wajib menyertakan foto",
            size=26, bold=True, color=GREEN_DEEP)
add_textbox(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
            "Berlaku untuk Critical, Preventif, dan Modifikasi — visual evidence langsung dari lapangan untuk verifikasi, audit, dan knowledge base.",
            size=12, color=INK_SOFT)

img_path = os.path.join(HERE, "powerops-photo-feature.png")
add_image_fit(s, img_path,
              Inches(0.6), Inches(2.0),
              Inches(12.1), Inches(3.8))

# 3 benefit points di bawah
benefits = [
    ("🔍", "Verifikasi cepat",  "Supervisor & foreman bisa konfirmasi kondisi tanpa harus turun ke lapangan."),
    ("📚", "Knowledge base",     "Foto tersimpan permanen → referensi untuk pekerjaan serupa di masa depan."),
    ("✅", "Audit-ready",        "Bukti visual + timestamp + operator → kepatuhan & traceability terjamin."),
]
y = Inches(5.95)
cw = Inches(4.0); gap = Inches(0.15)
x = Inches(0.6)
for icon, title, desc in benefits:
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, Inches(0.9))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = GREEN_LIGHT
    add_textbox(s, x + Inches(0.15), y + Inches(0.1), Inches(0.5), Inches(0.4),
                icon, size=18)
    add_textbox(s, x + Inches(0.7), y + Inches(0.06), cw - Inches(0.8), Inches(0.4),
                title, size=14, bold=True, color=INK)
    add_textbox(s, x + Inches(0.7), y + Inches(0.38), cw - Inches(0.8), Inches(0.5),
                desc, size=11, color=INK_SOFT)
    x += cw + gap

add_footer(s, 5, 6)
set_notes(s, """SPEAKER NOTES — Dokumentasi Foto

Penjelasan:
- Setiap kali operator/maintenance team membuat entry — apakah itu Critical, Preventif, atau Modifikasi — sistem mewajibkan upload foto.
- Foto ini menjadi bukti visual: kondisi sebelum tindakan, proses pengerjaan, dan hasil akhir (kalau perlu beberapa foto).
- Tiga tipe pekerjaan punya warna berbeda di UI (merah untuk Critical, biru untuk Preventif, oranye untuk Modifikasi) — supaya filter & sorting cepat.

Tiga manfaat utama (yang ada di kartu bawah slide):
1. Verifikasi cepat — supervisor tidak perlu turun ke lapangan untuk konfirmasi kondisi, cukup lihat foto + deskripsi di sistem.
2. Knowledge base — kalau 3 bulan lagi ada masalah serupa, kita tinggal cari foto pekerjaan sebelumnya untuk referensi cara penanganan.
3. Audit-ready — untuk kepatuhan internal & audit eksternal, kita punya bukti visual lengkap dengan timestamp, lokasi (kalau ada GPS), dan operator yang bertanggung jawab.

Talking point:
- "Selama ini banyak pekerjaan maintenance dilaporkan hanya dengan teks. 6 bulan kemudian kita lupa konteksnya. Dengan foto, ingatan organisasi terjaga."
- "Ini juga membantu transfer knowledge ke operator baru — mereka bisa belajar dari riwayat pekerjaan dengan visual yang jelas."

Antisipasi pertanyaan:
- "Apakah wajib foto saat entry, atau bisa diupload kemudian?" → bisa keduanya, tapi sistem akan mengingatkan kalau ada entry tanpa foto setelah 24 jam.
- "Bagaimana kalau foto besar / storage?" → otomatis di-compress saat upload, hemat storage tanpa kehilangan detail penting.""")

# =========================================================
# SLIDE 6 — Ringkasan
# =========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s)
add_textbox(s, Inches(0.4), Inches(0.08), Inches(10), Inches(0.4),
            "Ringkasan & Manfaat",
            size=20, bold=True, color=WHITE)

add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
            "Apa yang berubah, dan kenapa penting",
            size=28, bold=True, color=GREEN_DEEP)

# 4 KPI/highlight cards
cards = [
    ("−80%", "Tools manual",       "Dari 5 channel terpisah menjadi 1 platform PowerOps.", GREEN_DEEP),
    ("1×",   "Input data",         "Operator input sekali — sistem distribusi otomatis.", GREEN_MAIN),
    ("100%", "Pekerjaan terstruktur","Critical Job mengelompokkan maintenance secara kronologis.", AMBER),
    ("📸",   "Bukti visual",       "Foto wajib di setiap Critical / Preventif / Modifikasi.", BROWN),
]
y = Inches(2.0)
x = Inches(0.6)
cw = Inches(3.0); gap = Inches(0.13)
for big, label, desc, color in cards:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, Inches(2.5))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = GREEN_LIGHT
    # color strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.15))
    strip.fill.solid(); strip.fill.fore_color.rgb = color
    strip.line.fill.background()

    add_textbox(s, x + Inches(0.2), y + Inches(0.35), cw - Inches(0.4), Inches(1.0),
                big, size=42, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.2), y + Inches(1.4), cw - Inches(0.4), Inches(0.4),
                label, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.2), y + Inches(1.8), cw - Inches(0.4), Inches(0.7),
                desc, size=11, color=INK_SOFT, align=PP_ALIGN.CENTER)
    x += cw + gap

# big closing callout
y2 = Inches(4.9)
call = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y2,
                          Inches(12.1), Inches(1.8))
call.fill.solid(); call.fill.fore_color.rgb = GREEN_DEEP
call.line.fill.background()
add_textbox(s, Inches(0.9), y2 + Inches(0.2), Inches(11.5), Inches(0.5),
            "DAMPAK NYATA", size=12, bold=True, color=AMBER)
add_textbox(s, Inches(0.9), y2 + Inches(0.55), Inches(11.5), Inches(1.2),
            "Operator fokus mengoperasikan plant. Supervisor fokus mengambil keputusan. "
            "Manajemen melihat status pekerjaan utuh — dengan bukti foto — kapan saja.",
            size=18, color=WHITE)

add_footer(s, 6, 6)
set_notes(s, """SPEAKER NOTES — Ringkasan

Closing message:
- Tiga perbaikan ini saling melengkapi:
  1. Workflow laporan: efisiensi input & distribusi.
  2. Critical Job: struktur pekerjaan & kronologi.
  3. Foto: bukti visual & knowledge.

- Bersama-sama, PowerOps mengubah operasi UBB dari "banyak channel, banyak input, banyak salah" menjadi "satu platform, satu input, akuntabel & bisa diaudit".

Pesan untuk atasan:
- "Yang ingin saya tekankan: ini bukan tools baru yang menambah beban. Justru sebaliknya — kami memangkas pekerjaan administratif yang selama ini membebani operator dan foreman."
- "Investasi awal sudah berjalan. Yang dibutuhkan sekarang adalah dukungan untuk adopsi penuh — termasuk training singkat untuk operator senior."

Next step (kalau ditanya):
- Roll-out bertahap ke semua shift dalam 2 bulan.
- Review berkala (bulanan) untuk evaluasi adopsi & feedback dari user.
- Integrasi lanjutan: dashboard manajemen, alert otomatis untuk Critical Job yang tertunda, dll.

Tutup dengan:
- "Terima kasih. Mohon arahan dan masukan Bapak/Ibu."
- Siap menerima pertanyaan.""")

# ----- Save -----
out = os.path.join(HERE, "PowerOps-UBB-Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
